#!/usr/bin/env python3
"""핫스팟 제휴 제안서 PPT 생성기

캐치테이블 글로벌 X {핫스팟} 제안서를 자동 생성합니다.

사용법:
    python3 generate_proposal.py
    python3 generate_proposal.py --hotspot "올리브영 명동점" --region "명동"
    python3 generate_proposal.py --config hotspot_config.json
"""

from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
import argparse
import json
import os

# ============================================================
# 색상 & 스타일
# ============================================================
BRAND_RED = RGBColor(0xFF, 0x3D, 0x00)
TEXT_DARK = RGBColor(0x26, 0x26, 0x26)
TEXT_GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
LINE_GRAY = RGBColor(0xDD, 0xDD, 0xDD)

SLIDE_WIDTH = Cm(33.867)
SLIDE_HEIGHT = Cm(19.05)

# ============================================================
# 기본 설정
# ============================================================
DEFAULT_CONFIG = {
    "hotspot_name": "올리브영 명동점",
    "hotspot_short": "올리브영",
    "region": "명동",
    "campaign_concept": "K-Beauty × K-Food",
    "campaign_tagline": "뷰티 쇼핑과 미식을 하나의 여정으로",
}

CATCHTABLE_INFO = {
    "app_downloads": "120만+",
    "total_members": "180만+",
    "restaurants": "1.6만+",
    "languages": "20개 언어",
    "monthly_foreign_users": "50만+",
    "features": ["실시간 예약", "원격 웨이팅", "탐색 지도", "리뷰"],
}

# ============================================================
# 헬퍼 함수
# ============================================================

def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text,
                font_size=14, bold=False, color=TEXT_DARK,
                alignment=PP_ALIGN.LEFT, font_name="Arial",
                line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox


def add_section_header(slide, text):
    add_textbox(slide, Cm(2), Cm(1.2), Cm(15), Cm(0.8),
                text, font_size=14, bold=True, color=TEXT_DARK)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(2), Cm(2.2), Cm(29.87), Cm(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE_GRAY
    line.line.fill.background()


def add_slide_title(slide, text, top=Cm(2.8)):
    add_textbox(slide, Cm(2), top, Cm(29.87), Cm(2),
                text, font_size=24, bold=True, color=TEXT_DARK,
                alignment=PP_ALIGN.CENTER, line_spacing=1.3)


def add_subtitle(slide, text, top=Cm(5)):
    add_textbox(slide, Cm(2), top, Cm(29.87), Cm(1),
                text, font_size=14, color=TEXT_GRAY,
                alignment=PP_ALIGN.CENTER)


def add_badge(slide, left, top, width, height, label, color=BRAND_RED):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # 둥근 정도
    shape.adjustments[0] = 0.3

    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_top = Cm(0.1)
    tf.margin_bottom = Cm(0.1)
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Malgun Gothic"
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_content_box(slide, left, top, width, height, lines, font_size=12):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_LIGHT
    shape.line.fill.background()
    shape.adjustments[0] = 0.05

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.8)
    tf.margin_right = Cm(0.8)
    tf.margin_top = Cm(0.6)
    tf.margin_bottom = Cm(0.4)

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_DARK
        p.font.name = "Arial"
        p.space_after = Pt(5)
    return shape


def add_card(slide, left, top, width, height, value, label):
    """수치 카드 (캐치테이블 소개용)"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_LIGHT
    shape.line.fill.background()
    shape.adjustments[0] = 0.08

    # 상단 레드 악센트 라인
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left + Cm(0.8), top + Cm(0.5),
        Cm(2.5), Cm(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BRAND_RED
    accent.line.fill.background()

    # 수치
    add_textbox(slide, left, top + Cm(1.2), width, Cm(2),
                value, font_size=30, bold=True, color=BRAND_RED,
                alignment=PP_ALIGN.CENTER)
    # 라벨
    add_textbox(slide, left, top + Cm(3.5), width, Cm(1.5),
                label, font_size=12, color=TEXT_GRAY,
                alignment=PP_ALIGN.CENTER)


def add_arrow_between(slide, left, top, width=Cm(1.5), height=Cm(1.5)):
    """양방향 화살표 아이콘"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND_RED
    shape.line.fill.background()
    shape.adjustments[0] = 0.5  # 완전 원형

    tf = shape.text_frame
    tf.margin_top = Cm(0)
    tf.margin_bottom = Cm(0)
    p = tf.paragraphs[0]
    p.text = "↔"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


# ============================================================
# 슬라이드 1: 표지
# ============================================================
def create_cover(prs, cfg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_RED)

    # 좌측 상단 작은 텍스트
    add_textbox(slide, Cm(3), Cm(3), Cm(20), Cm(1),
                "Partnership Proposal", font_size=16, color=WHITE)

    # 메인 타이틀
    add_textbox(slide, Cm(3), Cm(5.5), Cm(27), Cm(4),
                f"캐치테이블 글로벌\nX {cfg['hotspot_name']}",
                font_size=44, bold=True, color=WHITE, line_spacing=1.3)

    # 캠페인 컨셉
    add_textbox(slide, Cm(3), Cm(10.5), Cm(27), Cm(1.5),
                f"<{cfg['campaign_concept']}> {cfg['campaign_tagline']}",
                font_size=20, color=WHITE)

    # 날짜
    add_textbox(slide, Cm(3), Cm(15.5), Cm(10), Cm(1),
                datetime.now().strftime("%Y.%m"), font_size=14, color=WHITE)


# ============================================================
# 슬라이드 2: AGENDA
# ============================================================
def create_agenda(prs, cfg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_RED)

    add_textbox(slide, Cm(3), Cm(3), Cm(27), Cm(2),
                "AGENDA", font_size=44, bold=True, color=WHITE)

    # 구분선
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(3), Cm(6), Cm(4), Cm(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = WHITE
    line.line.fill.background()

    items = [
        "01  캐치테이블 글로벌 소개",
        f"02  캐치테이블 글로벌 × {cfg['hotspot_short']} 캠페인 전략",
        "03  협업 방안 — 오프라인 연계 / 앱 내 노출",
    ]
    for i, item in enumerate(items):
        add_textbox(slide, Cm(3), Cm(7.5 + i * 2.2), Cm(27), Cm(1.5),
                    item, font_size=22, bold=True, color=WHITE)


# ============================================================
# 슬라이드 3: 캐치테이블 글로벌 소개
# ============================================================
def create_intro(prs, cfg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_section_header(slide, "캐치테이블 글로벌 소개")
    add_slide_title(slide, "외국인 관광객을 위한 No.1 맛집 예약 플랫폼")
    add_subtitle(slide,
        "한국을 방문하는 외국인 관광객에게 최고의 미식 경험을 연결합니다")

    # 5개 수치 카드
    info = CATCHTABLE_INFO
    metrics = [
        (info["app_downloads"], "글로벌 앱 다운로드"),
        (info["total_members"], "토탈 회원"),
        (info["restaurants"], "제휴 레스토랑"),
        (info["monthly_foreign_users"], "월간 외국인 유저"),
        (info["languages"], "지원 언어"),
    ]

    card_w = Cm(5.6)
    card_h = Cm(5.5)
    gap = Cm(0.5)
    start_x = Cm(1.5)
    y = Cm(6.5)

    for i, (value, label) in enumerate(metrics):
        x = start_x + i * (card_w + gap)
        add_card(slide, x, y, card_w, card_h, value, label)

    # 하단 주요 기능
    feat_text = "주요 기능:  " + "  ·  ".join(info["features"])
    add_textbox(slide, Cm(2), Cm(13.5), Cm(29.87), Cm(1),
                feat_text, font_size=13, color=TEXT_GRAY,
                alignment=PP_ALIGN.CENTER)

    # 하단 보충 설명
    add_textbox(slide, Cm(2), Cm(15), Cm(29.87), Cm(2),
                "캐치테이블은 한국 내 최대 레스토랑 예약 플랫폼으로,\n"
                "글로벌 버전을 통해 외국인 관광객 대상 서비스를 운영하고 있습니다.",
                font_size=11, color=TEXT_GRAY,
                alignment=PP_ALIGN.CENTER)


# ============================================================
# 슬라이드 4: 캠페인 전략
# ============================================================
def create_strategy(prs, cfg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_section_header(slide, "캠페인 전략")
    add_slide_title(slide,
        f"{cfg['hotspot_short']} 고객에게 근처 맛집 혜택을,\n"
        f"캐치테이블 유저에게 {cfg['hotspot_short']} 정보를")

    # ── 좌측: 핫스팟 혜택 ──
    add_badge(slide, Cm(2.5), Cm(6.2), Cm(6), Cm(1.1),
              f"{cfg['hotspot_short']}이 받는 혜택")

    hotspot_benefits = [
        f"캐치테이블 앱 배너에 {cfg['hotspot_short']} 홍보 노출",
        f"탐색 지도에 {cfg['hotspot_short']} 매장 위치 표기",
        "캐치테이블 SNS 채널을 통한 공동 홍보",
        f"월 {CATCHTABLE_INFO['monthly_foreign_users']} 외국인 유저에게 도달",
    ]
    add_content_box(slide, Cm(2.5), Cm(7.7), Cm(13), Cm(7),
                    [f"•  {b}" for b in hotspot_benefits], font_size=13)

    # ── 중앙 화살표 ──
    add_arrow_between(slide, Cm(16.2), Cm(10), Cm(1.5), Cm(1.5))

    # ── 우측: 캐치테이블 혜택 ──
    add_badge(slide, Cm(18.5), Cm(6.2), Cm(7), Cm(1.1),
              "캐치테이블이 받는 혜택")

    ct_benefits = [
        f"{cfg['hotspot_short']} 매장 내 캐치테이블 안내물 비치",
        f"{cfg['region']} 방문 고객의 맛집 예약 유도",
        f"{cfg['hotspot_short']} 고객 전용 할인 쿠폰 배포",
        "오프라인 접점을 통한 앱 신규 유입",
    ]
    add_content_box(slide, Cm(18.5), Cm(7.7), Cm(13), Cm(7),
                    [f"•  {b}" for b in ct_benefits], font_size=13)


# ============================================================
# 슬라이드 5: 오프라인 연계 — 안내물 비치
# ============================================================
def create_offline(prs, cfg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_section_header(slide, "협업 방안 — 오프라인 연계")
    add_slide_title(slide,
        f"{cfg['hotspot_short']} 매장 내 안내물 비치로\n"
        "고객에게 주변 맛집 할인 혜택 제공")

    # 프로세스
    add_badge(slide, Cm(2), Cm(6.2), Cm(3.8), Cm(1), "프로세스")
    process = [
        f"STEP 1.   {cfg['hotspot_short']} 매장 내 캐치테이블 안내물(포스터 / 테이블 스탠드 / 리플렛) 비치",
        f"STEP 2.   고객이 QR 코드를 스캔 → {cfg['region']} 주변 인기 맛집 리스트 + 전용 할인 쿠폰 자동 발급",
        "STEP 3.   캐치테이블 앱에서 바로 예약 또는 웨이팅 등록",
    ]
    add_content_box(slide, Cm(2), Cm(7.5), Cm(29.87), Cm(4.2),
                    process, font_size=13)

    # 기대효과
    add_badge(slide, Cm(2), Cm(12.5), Cm(3.8), Cm(1), "기대효과")
    effects = [
        f"•  {cfg['hotspot_short']} 고객에게 부가 혜택(맛집 할인)을 제공하여 매장 방문 만족도 향상",
        f"•  {cfg['region']} 인근 맛집으로의 자연스러운 동선 형성 → 지역 내 체류 시간 증가",
        "•  오프라인 접점에서 캐치테이블 앱 신규 유입 확보",
    ]
    add_content_box(slide, Cm(2), Cm(13.8), Cm(29.87), Cm(3.5),
                    effects, font_size=13)


# ============================================================
# 슬라이드 6: 앱 내 노출 — 배너 광고
# ============================================================
def create_banner(prs, cfg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_section_header(slide, "협업 방안 — 앱 내 노출")
    add_slide_title(slide,
        f"캐치테이블 앱 배너를 통한\n{cfg['hotspot_short']} 홍보 콘텐츠 노출")

    # 프로세스
    add_badge(slide, Cm(2), Cm(6.2), Cm(3.8), Cm(1), "프로세스")
    process = [
        f"STEP 1.   캐치테이블 앱 메인 배너 또는 인트로 팝업에 {cfg['hotspot_short']} 홍보 콘텐츠 게재",
        f"STEP 2.   배너 클릭 시 {cfg['hotspot_short']}이 원하는 프로모션 · 이벤트 · 매장 정보 페이지로 랜딩",
    ]
    add_content_box(slide, Cm(2), Cm(7.5), Cm(29.87), Cm(3.5),
                    process, font_size=13)

    # 기대효과
    add_badge(slide, Cm(2), Cm(11.8), Cm(3.8), Cm(1), "기대효과")
    effects = [
        f"•  캐치테이블 월간 외국인 유저 {CATCHTABLE_INFO['monthly_foreign_users']}에게 {cfg['hotspot_short']} 직접 노출",
        f"•  맛집을 탐색 중인 관광객에게 자연스러운 {cfg['hotspot_short']} 방문 동기 부여",
        "•  앱 내 가장 주목도 높은 지면(메인 배너)을 활용한 효율적 홍보",
        f"•  {cfg['hotspot_short']}이 원하는 시기 · 내용에 맞춰 유연하게 운영 가능",
    ]
    add_content_box(slide, Cm(2), Cm(13.1), Cm(29.87), Cm(4.2),
                    effects, font_size=13)


# ============================================================
# 슬라이드 7: 앱 내 노출 — 지도 & SNS
# ============================================================
def create_map_sns(prs, cfg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_section_header(slide, "협업 방안 — 앱 내 노출")
    add_slide_title(slide,
        f"탐색 지도 내 {cfg['hotspot_short']} 표기 + SNS 공동 홍보")

    # ── 좌측: 탐색 지도 ──
    add_badge(slide, Cm(2), Cm(5.8), Cm(4), Cm(1), "탐색 지도")

    map_process = [
        f"STEP 1.  캐치테이블 앱 탐색 지도에",
        f"             {cfg['hotspot_short']} 매장 위치를 아이콘으로 표기",
        f"STEP 2.  클릭 시 {cfg['hotspot_short']} 매장 정보 표시",
    ]
    add_content_box(slide, Cm(2), Cm(7.1), Cm(14.5), Cm(3.5),
                    map_process, font_size=12)

    map_effects = [
        f"•  맛집 검색 중 {cfg['hotspot_short']}을 자연스럽게 인지",
        f"•  식사 전후 {cfg['hotspot_short']} 방문 동선 형성",
        f"•  특히 {cfg['region']} 등 관광 핫스팟 내 방문객 증대",
    ]
    add_content_box(slide, Cm(2), Cm(11.2), Cm(14.5), Cm(3.5),
                    map_effects, font_size=12)

    # ── 우측: SNS 공동 홍보 ──
    add_badge(slide, Cm(17.5), Cm(5.8), Cm(5.5), Cm(1), "SNS 공동 홍보")

    sns_process = [
        "STEP 1.  캐치테이블 글로벌 SNS 채널에서",
        f"             {cfg['hotspot_short']} 관련 콘텐츠 공동 제작",
        "STEP 2.  Instagram · TikTok 등 다국어 게시",
    ]
    add_content_box(slide, Cm(17.5), Cm(7.1), Cm(14.37), Cm(3.5),
                    sns_process, font_size=12)

    sns_effects = [
        "•  양사 계정 태그를 통한 크로스 프로모션",
        f"•  {cfg['hotspot_short']} 관련 맛집 큐레이션 콘텐츠",
        "•  외국인 타겟 다국어 콘텐츠로 글로벌 도달",
    ]
    add_content_box(slide, Cm(17.5), Cm(11.2), Cm(14.37), Cm(3.5),
                    sns_effects, font_size=12)


# ============================================================
# 슬라이드 8: 마무리
# ============================================================
def create_closing(prs, cfg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_RED)

    add_textbox(slide, Cm(2), Cm(6.5), Cm(29.87), Cm(3),
                "END OF DOCUMENT", font_size=54, bold=True,
                color=WHITE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Cm(2), Cm(10.5), Cm(29.87), Cm(1.5),
                "캐치테이블 글로벌팀  |  catchtable.global",
                font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)


# ============================================================
# 메인
# ============================================================
def generate_proposal(cfg=None):
    if cfg is None:
        cfg = DEFAULT_CONFIG

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    create_cover(prs, cfg)
    create_agenda(prs, cfg)
    create_intro(prs, cfg)
    create_strategy(prs, cfg)
    create_offline(prs, cfg)
    create_banner(prs, cfg)
    create_map_sns(prs, cfg)
    create_closing(prs, cfg)

    os.makedirs("output", exist_ok=True)
    safe_name = cfg['hotspot_name'].replace(' ', '_').replace('/', '_')
    filename = f"output/[캐치테이블 글로벌] {safe_name} 제휴 제안서.pptx"
    prs.save(filename)
    print(f"✅ 제안서 생성 완료: {filename}")
    print(f"   슬라이드 {len(prs.slides)}장")
    return filename


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="핫스팟 제휴 제안서 생성기")
    parser.add_argument("--config", type=str, help="JSON 설정 파일 경로")
    parser.add_argument("--hotspot", type=str, help="핫스팟 이름 (전체)")
    parser.add_argument("--short", type=str, help="핫스팟 약칭")
    parser.add_argument("--region", type=str, help="지역명")
    parser.add_argument("--concept", type=str, help="캠페인 컨셉")
    parser.add_argument("--tagline", type=str, help="캠페인 태그라인")
    args = parser.parse_args()

    cfg = DEFAULT_CONFIG.copy()

    if args.config:
        with open(args.config, 'r', encoding='utf-8') as f:
            cfg.update(json.load(f))

    if args.hotspot:
        cfg['hotspot_name'] = args.hotspot
    if args.short:
        cfg['hotspot_short'] = args.short
    elif args.hotspot and 'hotspot_short' not in cfg:
        cfg['hotspot_short'] = args.hotspot
    if args.region:
        cfg['region'] = args.region
    if args.concept:
        cfg['campaign_concept'] = args.concept
    if args.tagline:
        cfg['campaign_tagline'] = args.tagline

    generate_proposal(cfg)
